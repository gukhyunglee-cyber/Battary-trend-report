"""
PPT Generator Module
Creates professional PowerPoint presentations from structured slide data.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime


class PPTGenerator:
    # Color palette
    DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)     # Dark navy
    ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)  # Professional blue
    ACCENT_GREEN = RGBColor(0x00, 0xA8, 0x6B) # Battery green
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
    DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
    MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
    SOURCE_COLOR = RGBColor(0x00, 0x66, 0xCC)

    SLIDE_COLORS = [
        RGBColor(0x00, 0x7A, 0xCC),  # Blue
        RGBColor(0x00, 0xA8, 0x6B),  # Green
        RGBColor(0xE6, 0x7E, 0x22),  # Orange
        RGBColor(0x8E, 0x44, 0xAD),  # Purple
    ]

    def __init__(self):
        self.prs = Presentation()
        # Set slide dimensions (widescreen 16:9)
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _set_slide_bg(self, slide, color):
        """Set solid background color for a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=14, font_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                     font_name="맑은 고딕"):
        """Helper to add a styled text box."""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                          Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        if font_color:
            p.font.color.rgb = font_color
        return tf

    def _create_title_slide(self, title_text="2차전지 업계 동향 리포트"):
        """Create a professional title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._set_slide_bg(slide, self.DARK_BG)

        today = datetime.now().strftime("%Y년 %m월 %d일")

        # Main title
        self._add_textbox(slide, 1.5, 1.8, 10, 1.5,
                          "⚡ " + title_text,
                          font_size=36, font_color=self.WHITE, bold=True,
                          alignment=PP_ALIGN.CENTER)

        # Subtitle
        self._add_textbox(slide, 2, 3.5, 9, 0.8,
                          "Battery Industry Weekly Trend Report",
                          font_size=20, font_color=self.ACCENT_BLUE,
                          alignment=PP_ALIGN.CENTER)

        # Date
        self._add_textbox(slide, 2, 4.5, 9, 0.6,
                          f"작성일: {today}  |  AI 분석 리포트 (Powered by Gemini)",
                          font_size=14, font_color=self.MEDIUM_GRAY,
                          alignment=PP_ALIGN.CENTER)

        # Bottom line
        line = slide.shapes.add_shape(
            1, Inches(2), Inches(5.5), Inches(9), Emu(28000)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.ACCENT_BLUE
        line.line.fill.background()

    def _create_content_slide(self, slide_info, slide_index=0):
        """Create a content slide with colored accent bar."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        self._set_slide_bg(slide, self.WHITE)

        accent_color = self.SLIDE_COLORS[slide_index % len(self.SLIDE_COLORS)]
        title = slide_info.get('title', 'N/A')
        content_lines = slide_info.get('content', [])

        # Issue number badge
        badge_text = f"Issue {slide_index + 1}"
        self._add_textbox(slide, 0.8, 0.4, 1.5, 0.4,
                          badge_text,
                          font_size=12, font_color=accent_color, bold=True)

        # Accent bar (left side)
        bar = slide.shapes.add_shape(
            1, Inches(0.5), Inches(0.9), Emu(60000), Inches(0.6)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

        # Title
        self._add_textbox(slide, 0.9, 0.85, 11, 0.7,
                          title,
                          font_size=26, font_color=self.DARK_TEXT, bold=True)

        # Content bullets
        if isinstance(content_lines, str):
            content_lines = [content_lines]

        y_pos = 1.8
        for line in content_lines:
            if not line or line.strip() == "":
                y_pos += 0.15
                continue

            is_source = line.strip().startswith("📎") or line.strip().startswith("http")
            is_url = "http" in line

            font_clr = self.SOURCE_COLOR if is_url else self.DARK_TEXT
            fsize = 11 if is_source else 14

            if is_source or is_url:
                prefix = ""
            else:
                prefix = "• " if not line.startswith("  ") else "  "

            self._add_textbox(slide, 1.0, y_pos, 11, 0.35,
                              prefix + line.strip(),
                              font_size=fsize, font_color=font_clr)
            y_pos += 0.38

        # Bottom accent line
        bottom_line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(6.9), Inches(12), Emu(18000)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = accent_color
        bottom_line.line.fill.background()

    def _create_summary_slide(self):
        """Create a closing summary slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_slide_bg(slide, self.DARK_BG)

        self._add_textbox(slide, 2, 2.5, 9, 1,
                          "Thank You",
                          font_size=40, font_color=self.WHITE, bold=True,
                          alignment=PP_ALIGN.CENTER)

        self._add_textbox(slide, 2, 3.8, 9, 0.6,
                          "본 리포트는 AI가 자동 수집·분석한 데이터를 기반으로 작성되었습니다.",
                          font_size=14, font_color=self.MEDIUM_GRAY,
                          alignment=PP_ALIGN.CENTER)

        today = datetime.now().strftime("%Y.%m.%d")
        self._add_textbox(slide, 2, 4.6, 9, 0.6,
                          f"Battery Trend Reporter  |  {today}",
                          font_size=12, font_color=self.ACCENT_BLUE,
                          alignment=PP_ALIGN.CENTER)

    def create_presentation(self, slides_data, output_file="battery_trend_report.pptx"):
        """
        Creates a full PPT presentation from structured slide data.
        slides_data: list of dicts -> [{'title': '...', 'content': ['...']}]
        """
        # Title slide
        self._create_title_slide()

        # Content slides
        for i, slide_info in enumerate(slides_data):
            self._create_content_slide(slide_info, i)

        # Summary slide
        self._create_summary_slide()

        # Save
        self.prs.save(output_file)
        print(f"[PPT Generator] Saved presentation to {output_file}")
        return os.path.abspath(output_file)


if __name__ == "__main__":
    gen = PPTGenerator()
    data = [
        {
            'title': '전고체 배터리 양산 경쟁 가속화',
            'content': [
                '삼성SDI, 2026년 전고체 배터리 파일럿 라인 가동 시작',
                '토요타, 2027년 전고체 배터리 탑재 EV 출시 계획 발표',
                '에너지 밀도 500Wh/kg 이상으로 기존 리튬이온 대비 2배',
                '주요 과제: 대량생산 비용 절감 및 계면 저항 문제 해결',
                '시장 규모 2030년까지 150억 달러 전망',
                '',
                '📎 출처:',
                'https://example.com/solid-state'
            ]
        },
        {
            'title': '나트륨이온 배터리 상용화 진전',
            'content': [
                'CATL, 나트륨이온 2세대 셀 양산 개시',
                'kWh당 40달러 이하 원가 달성으로 가격 혁신',
                'ESS 및 저가형 EV 시장에서 LFP 대체 가능성',
                '리튬 의존도 감소로 공급망 안정성 확보',
                '한계: 에너지 밀도 160Wh/kg 수준',
            ]
        },
    ]
    gen.create_presentation(data, "test_professional.pptx")
    print("Test PPT generated!")
